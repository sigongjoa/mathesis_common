import logging
import json
from typing import List, Dict, Any, Optional
from pptx import Presentation
from pptx.util import Inches, Pt

from mathesis_core.llm.clients import OllamaClient

logger = logging.getLogger(__name__)

class PPTGeneratorAgent:
    """
    Agent that generates PowerPoint presentations using LLM for content 
    and python-pptx for rendering.
    """

    def __init__(self, ollama_client: Optional[OllamaClient] = None):
        self.llm = ollama_client or OllamaClient()

    async def generate_presentation(self, topic: str, num_slides: int = 5, output_path: str = "output.pptx") -> str:
        """
        Generates a PPTX file based on a topic.
        
        Args:
            topic: The subject of the presentation.
            num_slides: Approximate number of slides.
            output_path: File path to save the PPTX.
            
        Returns:
            Absolute path to the generated PPTX.
        """
        logger.info(f"Generating PPT for topic: '{topic}' with {num_slides} slides")

        # 1. Generate Content Structure via LLM
        structure = await self._generate_structure(topic, num_slides)
        
        # 2. Render to PPTX
        self._render_pptx(structure, output_path)
        
        logger.info(f"PPT generated successfully at {output_path}")
        return output_path

    async def _generate_structure(self, topic: str, num_slides: int) -> List[Dict[str, Any]]:
        """
        Uses LLM to create a JSON structure for the slides.
        """
        prompt = f"""
        You are a presentation expert. Create a structured outline for a PowerPoint presentation about "{topic}".
        The presentation should have exactly {num_slides} slides.
        
        Return ONLY valid JSON in the following format:
        [
            {{
                "title": "Slide Title",
                "content": ["Bullet point 1", "Bullet point 2", "Bullet point 3"],
                "notes": "Speaker notes for this slide"
            }},
            ...
        ]
        
        Do not include markdown formatting like ```json. Just raw JSON.
        """
        
        response_text = await self.llm.generate_text(prompt)
        
        # Clean up potential markdown code blocks
        clean_text = response_text.replace("```json", "").replace("```", "").strip()
        
        try:
            slides = json.loads(clean_text)
            return slides
        except json.JSONDecodeError:
            logger.error(f"Failed to parse LLM response as JSON: {response_text}")
            # Fallback for resilience
            return [
                {
                    "title": f"Presentation on {topic}",
                    "content": ["Failed to generate structured content.", "Please try again."],
                    "notes": "Error in generation."
                }
            ]

    def _render_pptx(self, slides_data: List[Dict[str, Any]], output_path: str):
        """
        Renders the JSON slide data into a .pptx file.
        """
        prs = Presentation()

        for slide_data in slides_data:
            # layout 1 is usually 'Title and Content'
            slide_layout = prs.slide_layouts[1] 
            slide = prs.slides.add_slide(slide_layout)

            # Set Title
            title = slide.shapes.title
            if title:
                title.text = slide_data.get("title", "Untitled")

            # Set Content (Bullet points)
            body_shape = slide.placeholders[1]
            tf = body_shape.text_frame
            tf.text = "" # Clear existing

            content_lines = slide_data.get("content", [])
            if isinstance(content_lines, str):
                content_lines = [content_lines]

            for i, line in enumerate(content_lines):
                if i == 0:
                    p = tf.paragraphs[0]
                else:
                    p = tf.add_paragraph()
                p.text = line
                p.level = 0

            # Set Notes
            if "notes" in slide_data:
                notes_slide = slide.notes_slide
                text_frame = notes_slide.notes_text_frame
                text_frame.text = slide_data["notes"]

        prs.save(output_path)
